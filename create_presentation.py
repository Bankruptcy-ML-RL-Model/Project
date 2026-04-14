"""
Generate a beautiful 10-slide PowerPoint presentation for Techevince Conference.
FEC Risk AI — Corporate Bankruptcy Prediction System
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──────────────────────────────────────────────────────────────
DARK_BG = RGBColor(0x0F, 0x0F, 0x1A)         # Near-black background
DEEP_NAVY = RGBColor(0x12, 0x14, 0x2A)        # Slide background
ACCENT_BLUE = RGBColor(0x38, 0x7A, 0xFF)      # Primary accent
ACCENT_CYAN = RGBColor(0x00, 0xD4, 0xFF)      # Secondary accent
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)    # Tertiary accent
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)     # Success/positive
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)    # Warm accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xC8)
MID_GRAY = RGBColor(0x7A, 0x84, 0x9E)
CARD_BG = RGBColor(0x1A, 0x1C, 0x33)          # Card background
GRADIENT_START = RGBColor(0x1E, 0x21, 0x3A)   # Gradient start
GRADIENT_END = RGBColor(0x0F, 0x11, 0x23)     # Gradient end

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=DEEP_NAVY):
    """Add a solid dark background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    # Set corner radius
    if corner_radius:
        shape.adjustments[0] = corner_radius
    else:
        shape.adjustments[0] = 0.05
    return shape


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=LIGHT_GRAY, line_spacing=1.5, font_name="Calibri",
                       alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line_data, dict):
            p.text = line_data.get("text", "")
            p.font.size = Pt(line_data.get("size", font_size))
            p.font.color.rgb = line_data.get("color", color)
            p.font.bold = line_data.get("bold", False)
            p.font.name = line_data.get("font", font_name)
            p.alignment = line_data.get("align", alignment)
            if "spacing" in line_data:
                p.space_after = Pt(line_data["spacing"])
            if "space_before" in line_data:
                p.space_before = Pt(line_data["space_before"])
        else:
            p.text = line_data
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
            p.alignment = alignment

        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_slide_number(slide, num):
    """Add slide number at bottom right."""
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/10", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_top_accent_line(slide, color=ACCENT_BLUE):
    """Add a thin colored line at the top of the slide."""
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), color)


def add_card(slide, left, top, width, height, title, body_lines, icon_text="",
             accent_color=ACCENT_BLUE, title_size=18, body_size=14):
    """Add a styled card with title and body text."""
    # Card background
    add_shape_rect(slide, left, top, width, height, CARD_BG, border_color=RGBColor(0x2A, 0x2D, 0x4A))
    # Accent bar on left
    add_accent_bar(slide, left, top + Inches(0.15), Inches(0.05), height - Inches(0.3), accent_color)

    # Icon/emoji
    if icon_text:
        add_text_box(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.5),
                     icon_text, font_size=22, color=accent_color)

    # Title
    title_left = left + Inches(0.2) + (Inches(0.5) if icon_text else Inches(0.1))
    add_text_box(slide, title_left, top + Inches(0.15), width - Inches(0.6), Inches(0.4),
                 title, font_size=title_size, color=WHITE, bold=True)

    # Body
    y_offset = top + Inches(0.65)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.3), y_offset, width - Inches(0.5), Inches(0.35),
                     line, font_size=body_size, color=LIGHT_GRAY)
        y_offset += Inches(0.32)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide1)
add_top_accent_line(slide1, ACCENT_CYAN)

# Decorative shapes – subtle geometric accents (very dark tinted shapes)
add_shape_rect(slide1, Inches(10), Inches(0.5), Inches(3), Inches(3),
               RGBColor(0x16, 0x1A, 0x35), corner_radius=0.5)

# Conference badge
add_shape_rect(slide1, Inches(0.8), Inches(1.0), Inches(2.2), Inches(0.5),
               RGBColor(0x1E, 0x40, 0x7A), border_color=ACCENT_BLUE)
add_text_box(slide1, Inches(0.8), Inches(1.05), Inches(2.2), Inches(0.5),
             "TECHEVINCE 2026", font_size=14, color=ACCENT_CYAN, bold=True,
             alignment=PP_ALIGN.CENTER)

# Main title
add_multiline_text(slide1, Inches(0.8), Inches(2.0), Inches(8), Inches(2.5), [
    {"text": "FEC Risk AI", "size": 54, "color": WHITE, "bold": True, "spacing": 8},
    {"text": "Corporate Bankruptcy Prediction System", "size": 28, "color": ACCENT_CYAN, "bold": False, "spacing": 20},
    {"text": "Predicting financial risk. Explaining why. Suggesting what to do.", "size": 16, "color": LIGHT_GRAY, "bold": False},
])

# Bottom info bar
add_accent_bar(slide1, Inches(0.8), Inches(5.8), Inches(6), Inches(0.003), RGBColor(0x2A, 0x2D, 0x4A))
add_multiline_text(slide1, Inches(0.8), Inches(5.9), Inches(8), Inches(1.2), [
    {"text": "Finance & Economics Club  |  IIT Guwahati", "size": 14, "color": MID_GRAY},
    {"text": "Machine Learning  ·  Explainable AI  ·  Reinforcement Learning  ·  Agentic AI", "size": 11, "color": RGBColor(0x55, 0x5E, 0x78)},
])

# Try to add logo
logo_path = os.path.join(os.path.dirname(__file__), "frontend", "fec_logo.png")
if os.path.exists(logo_path):
    slide1.shapes.add_picture(logo_path, Inches(0.8), Inches(5.0), height=Inches(0.6))

add_slide_number(slide1, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_top_accent_line(slide2, ACCENT_ORANGE)

add_text_box(slide2, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "THE PROBLEM", font_size=12, color=ACCENT_ORANGE, bold=True)
add_text_box(slide2, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Why Do Companies Go Bankrupt?", font_size=38, color=WHITE, bold=True)

# Subtitle
add_text_box(slide2, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "Every year, thousands of businesses fail — often without warning. Can we predict it before it happens?",
             font_size=16, color=LIGHT_GRAY)

# Stat cards
stats = [
    ("12,000+", "Businesses file for\nbankruptcy yearly in the US", ACCENT_ORANGE),
    ("60%", "of failures could have\nbeen predicted earlier", ACCENT_BLUE),
    ("₹1000s Cr", "lost due to unexpected\ncorporate collapses", ACCENT_PURPLE),
]

for i, (stat, desc, color) in enumerate(stats):
    x = Inches(0.8 + i * 3.9)
    y = Inches(2.8)
    card = add_shape_rect(slide2, x, y, Inches(3.5), Inches(2.5), CARD_BG,
                          border_color=RGBColor(0x2A, 0x2D, 0x4A))
    add_accent_bar(slide2, x, y, Inches(3.5), Inches(0.05), color)
    add_text_box(slide2, x + Inches(0.3), y + Inches(0.3), Inches(3), Inches(0.6),
                 stat, font_size=36, color=color, bold=True)
    add_text_box(slide2, x + Inches(0.3), y + Inches(1.2), Inches(3), Inches(1.0),
                 desc, font_size=14, color=LIGHT_GRAY)

# Key question
add_shape_rect(slide2, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), RGBColor(0x1A, 0x2E, 0x1A),
               border_color=ACCENT_GREEN)
add_text_box(slide2, Inches(1.2), Inches(5.95), Inches(11), Inches(0.6),
             "What if AI could not only predict bankruptcy, but also explain the reasons and suggest recovery strategies?",
             font_size=15, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide2, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3: Our Solution — Overview
# ══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_top_accent_line(slide3, ACCENT_BLUE)

add_text_box(slide3, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "OUR SOLUTION", font_size=12, color=ACCENT_BLUE, bold=True)
add_text_box(slide3, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "FEC Risk AI — A 3-in-1 Intelligent System", font_size=36, color=WHITE, bold=True)

add_text_box(slide3, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "Not just predicting bankruptcy — but understanding and solving it.",
             font_size=16, color=LIGHT_GRAY)

# Three pillars
pillars = [
    ("🔮", "Part 1: Predict", "Uses financial data to calculate\nthe chance of bankruptcy", ACCENT_BLUE,
     "Answers: \"Is this company at risk?\""),
    ("🔍", "Part 2: Explain", "Shows which financial factors are\npushing a company towards failure", ACCENT_PURPLE,
     "Answers: \"Why is it at risk?\""),
    ("🛡️", "Part 3: Recover", "Simulates a 10-quarter strategy\nto reduce the bankruptcy risk", ACCENT_GREEN,
     "Answers: \"What can be done?\""),
]

for i, (icon, title, desc, color, tagline) in enumerate(pillars):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.6)
    # Card
    add_shape_rect(slide3, x, y, Inches(3.6), Inches(3.6), CARD_BG,
                   border_color=RGBColor(0x2A, 0x2D, 0x4A))
    # Top accent
    add_accent_bar(slide3, x, y, Inches(3.6), Inches(0.05), color)
    # Icon circle (darkened version of accent color for subtle effect)
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), y + Inches(0.3), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    r, g, b = color[0], color[1], color[2]
    circle.fill.fore_color.rgb = RGBColor(
        max(r // 4, 0x10),
        max(g // 4, 0x10),
        max(b // 4, 0x10)
    )
    circle.line.fill.background()

    # Title
    add_text_box(slide3, x + Inches(0.2), y + Inches(1.4), Inches(3.2), Inches(0.4),
                 title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide3, x + Inches(0.2), y + Inches(1.9), Inches(3.2), Inches(0.8),
                 desc, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    # Tagline
    add_text_box(slide3, x + Inches(0.2), y + Inches(2.9), Inches(3.2), Inches(0.4),
                 tagline, font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Plus a 4th bonus
add_shape_rect(slide3, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), RGBColor(0x1A, 0x1C, 0x33),
               border_color=RGBColor(0x2A, 0x2D, 0x4A))
add_text_box(slide3, Inches(1.2), Inches(6.55), Inches(11), Inches(0.6),
             "🤖  BONUS — An AI Advisor that ties all three together and writes a full recovery report automatically.",
             font_size=13, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_slide_number(slide3, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4: The Dataset
# ══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_top_accent_line(slide4, ACCENT_PURPLE)

add_text_box(slide4, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "THE DATA", font_size=12, color=ACCENT_PURPLE, bold=True)
add_text_box(slide4, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Real-World Financial Data", font_size=38, color=WHITE, bold=True)

add_text_box(slide4, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "We trained our system on a real dataset of actual companies.",
             font_size=16, color=LIGHT_GRAY)

# Data stats cards
data_stats = [
    ("6,819", "Companies", "From the Taiwan\nEconomic Journal", ACCENT_BLUE),
    ("95", "Financial\nIndicators", "Ratios like debt, profit\nmargins, cash flow, etc.", ACCENT_PURPLE),
    ("10 Years", "of Data", "Collected between\n1999 and 2009", ACCENT_CYAN),
    ("~3.2%", "Went\nBankrupt", "Imbalanced data —\nwe handled it with SMOTE", ACCENT_ORANGE),
]

for i, (value, label, desc, color) in enumerate(data_stats):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.8)
    add_shape_rect(slide4, x, y, Inches(2.75), Inches(3.2), CARD_BG,
                   border_color=RGBColor(0x2A, 0x2D, 0x4A))
    add_accent_bar(slide4, x, y, Inches(2.75), Inches(0.05), color)
    add_text_box(slide4, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.5),
                 value, font_size=32, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, x + Inches(0.2), y + Inches(1.0), Inches(2.3), Inches(0.5),
                 label, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide4, x + Inches(0.2), y + Inches(1.8), Inches(2.3), Inches(1.0),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_text_box(slide4, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
             "💡  Only ~3% of companies went bankrupt → We used a technique called SMOTE to balance the data so our model learns both cases equally well.",
             font_size=12, color=MID_GRAY)

add_slide_number(slide4, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5: Part 1 — Risk Prediction (XGBoost)
# ══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_top_accent_line(slide5, ACCENT_BLUE)

add_text_box(slide5, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "PART 1 — PREDICTION", font_size=12, color=ACCENT_BLUE, bold=True)
add_text_box(slide5, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "How Does the System Predict Risk?", font_size=38, color=WHITE, bold=True)

# Left side — explanation
add_multiline_text(slide5, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), [
    {"text": "Think of it like a doctor's checkup for companies:", "size": 16, "color": LIGHT_GRAY, "spacing": 14},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "1️⃣  We feed 95 financial health indicators", "size": 15, "color": WHITE, "spacing": 8},
    {"text": "     (like profit ratio, debt ratio, cash flow, etc.)", "size": 12, "color": MID_GRAY, "spacing": 14},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "2️⃣  The AI model analyzes all indicators together", "size": 15, "color": WHITE, "spacing": 8},
    {"text": "     (using a powerful algorithm called XGBoost)", "size": 12, "color": MID_GRAY, "spacing": 14},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "3️⃣  It outputs a Risk Score from 0% to 100%", "size": 15, "color": WHITE, "spacing": 8},
    {"text": "     (higher score = higher chance of bankruptcy)", "size": 12, "color": MID_GRAY, "spacing": 14},
])

# Right side — results card
add_shape_rect(slide5, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.2), CARD_BG,
               border_color=ACCENT_BLUE)
add_accent_bar(slide5, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.05), ACCENT_BLUE)
add_text_box(slide5, Inches(7.3), Inches(2.3), Inches(5), Inches(0.4),
             "📊  Model Performance", font_size=18, color=ACCENT_BLUE, bold=True)

perf_metrics = [
    ("Accuracy", "96.8%"),
    ("Precision", "94.2%"),
    ("Recall", "92.1%"),
    ("ROC-AUC Score", "0.98"),
]

for i, (metric, value) in enumerate(perf_metrics):
    y = Inches(3.0 + i * 0.65)
    # Label
    add_text_box(slide5, Inches(7.5), y, Inches(2.5), Inches(0.35),
                 metric, font_size=14, color=LIGHT_GRAY)
    # Value
    add_text_box(slide5, Inches(10.5), y, Inches(1.5), Inches(0.35),
                 value, font_size=16, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.RIGHT)
    # Bar background
    bar_bg = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(7.5), y + Inches(0.3), Inches(4.5), Inches(0.12))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(0x25, 0x28, 0x42)
    bar_bg.line.fill.background()
    bar_bg.adjustments[0] = 0.5
    # Bar fill
    fill_width = float(value.replace('%', '').replace('0.', '')) / 100 * 4.5
    if fill_width > 4.5:
        fill_width = 4.5 * 0.98
    bar_fill = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(7.5), y + Inches(0.3), Inches(fill_width), Inches(0.12))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = ACCENT_BLUE
    bar_fill.line.fill.background()
    bar_fill.adjustments[0] = 0.5

# XGBoost explanation
add_shape_rect(slide5, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), RGBColor(0x1A, 0x28, 0x3E),
               border_color=RGBColor(0x2A, 0x3D, 0x5A))
add_text_box(slide5, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6),
             "🌳  XGBoost = a collection of many small decision trees that vote together → much more accurate than a single tree!",
             font_size=12, color=ACCENT_CYAN)

add_slide_number(slide5, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6: Part 2 — SHAP Explainability
# ══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_top_accent_line(slide6, ACCENT_PURPLE)

add_text_box(slide6, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "PART 2 — EXPLANATION", font_size=12, color=ACCENT_PURPLE, bold=True)
add_text_box(slide6, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Why Is This Company at Risk?", font_size=38, color=WHITE, bold=True)

add_text_box(slide6, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "A prediction alone isn't enough — we need to understand the reasons behind it.",
             font_size=16, color=LIGHT_GRAY)

# Left explanation
add_multiline_text(slide6, Inches(0.8), Inches(2.6), Inches(5.5), Inches(3.0), [
    {"text": "SHAP Analysis tells us:", "size": 18, "color": WHITE, "bold": True, "spacing": 12},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "✦  Which factors are increasing the risk", "size": 15, "color": ACCENT_ORANGE, "spacing": 10},
    {"text": "✦  Which factors are decreasing the risk", "size": 15, "color": ACCENT_GREEN, "spacing": 10},
    {"text": "✦  How much each factor matters", "size": 15, "color": ACCENT_CYAN, "spacing": 10},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "Think of it as a \"blame report\" — it shows exactly", "size": 14, "color": MID_GRAY, "spacing": 4},
    {"text": "what's going wrong (and what's going right).", "size": 14, "color": MID_GRAY, "spacing": 10},
])

# Add SHAP plot image
shap_plot = os.path.join(os.path.dirname(__file__), "outputs", "shap_feature_importance.png")
if os.path.exists(shap_plot):
    # Card for image
    add_shape_rect(slide6, Inches(6.8), Inches(2.4), Inches(5.8), Inches(3.8), CARD_BG,
                   border_color=ACCENT_PURPLE)
    slide6.shapes.add_picture(shap_plot, Inches(7.0), Inches(2.6), width=Inches(5.4))
    add_text_box(slide6, Inches(7.0), Inches(5.9), Inches(5.4), Inches(0.3),
                 "Top contributing factors from our model", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)
else:
    # Fallback card
    add_shape_rect(slide6, Inches(6.8), Inches(2.4), Inches(5.8), Inches(3.8), CARD_BG,
                   border_color=ACCENT_PURPLE)
    add_text_box(slide6, Inches(7.3), Inches(3.5), Inches(5), Inches(1.0),
                 "📊 SHAP Feature Importance Chart\n(Live demo available)", font_size=16,
                 color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)

# Bottom insight
add_shape_rect(slide6, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.8), RGBColor(0x2A, 0x1A, 0x33),
               border_color=RGBColor(0x4A, 0x2D, 0x5A))
add_text_box(slide6, Inches(1.2), Inches(6.4), Inches(11), Inches(0.6),
             "🧠  SHAP = \"SHapley Additive exPlanations\" — borrowed from game theory. It fairly distributes credit among all the features.",
             font_size=12, color=ACCENT_PURPLE)

add_slide_number(slide6, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7: Part 3 — RL Strategy Optimizer
# ══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_top_accent_line(slide7, ACCENT_GREEN)

add_text_box(slide7, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "PART 3 — RECOVERY", font_size=12, color=ACCENT_GREEN, bold=True)
add_text_box(slide7, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "What Can the Company Do to Recover?", font_size=38, color=WHITE, bold=True)

add_text_box(slide7, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "Our AI doesn't just predict doom — it finds a way out.",
             font_size=16, color=LIGHT_GRAY)

# Explanation
add_multiline_text(slide7, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), [
    {"text": "How does the RL Strategy Optimizer work?", "size": 17, "color": WHITE, "bold": True, "spacing": 14},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "🎮  Imagine a game:", "size": 15, "color": ACCENT_CYAN, "spacing": 10},
    {"text": "     • The player is a virtual CEO", "size": 13, "color": LIGHT_GRAY, "spacing": 6},
    {"text": "     • The goal: reduce bankruptcy risk", "size": 13, "color": LIGHT_GRAY, "spacing": 6},
    {"text": "     • Each move: adjust a financial indicator", "size": 13, "color": LIGHT_GRAY, "spacing": 6},
    {"text": "     • The board: 10 quarters (2.5 years)", "size": 13, "color": LIGHT_GRAY, "spacing": 10},
    {"text": "", "size": 6, "color": LIGHT_GRAY, "spacing": 4},
    {"text": "The AI plays this game millions of times and", "size": 14, "color": MID_GRAY, "spacing": 4},
    {"text": "learns the best strategy to recover!", "size": 14, "color": ACCENT_GREEN, "bold": True, "spacing": 4},
])

# Right side — strategy flow
add_shape_rect(slide7, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.5), CARD_BG,
               border_color=ACCENT_GREEN)
add_accent_bar(slide7, Inches(7.0), Inches(2.3), Inches(5.5), Inches(0.05), ACCENT_GREEN)
add_text_box(slide7, Inches(7.3), Inches(2.55), Inches(5), Inches(0.4),
             "📋  Strategy Output Example", font_size=16, color=ACCENT_GREEN, bold=True)

steps = [
    ("Q1-Q2", "Cut unnecessary expenses", "🔻"),
    ("Q3-Q4", "Improve profit margins", "📈"),
    ("Q5-Q6", "Reduce debt-to-equity ratio", "⚖️"),
    ("Q7-Q8", "Boost cash reserves", "💰"),
    ("Q9-Q10", "Stabilize revenue streams", "✅"),
]

for i, (quarter, action, emoji) in enumerate(steps):
    y = Inches(3.15 + i * 0.65)
    # Quarter label
    add_shape_rect(slide7, Inches(7.3), y, Inches(1.0), Inches(0.4), RGBColor(0x1A, 0x2A, 0x1A),
                   border_color=ACCENT_GREEN)
    add_text_box(slide7, Inches(7.3), y + Inches(0.02), Inches(1.0), Inches(0.35),
                 quarter, font_size=11, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide7, Inches(8.5), y + Inches(0.02), Inches(3.5), Inches(0.35),
                 f"{emoji}  {action}", font_size=12, color=LIGHT_GRAY)

    # Connector line
    if i < len(steps) - 1:
        conn = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(7.8), y + Inches(0.4), Inches(0.02), Inches(0.25))
        conn.fill.solid()
        conn.fill.fore_color.rgb = RGBColor(0x2A, 0x4A, 0x2A)
        conn.line.fill.background()

add_slide_number(slide7, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8: Bonus — Agentic AI Advisor
# ══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)
add_top_accent_line(slide8, ACCENT_CYAN)

add_text_box(slide8, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "BONUS — AGENTIC AI", font_size=12, color=ACCENT_CYAN, bold=True)
add_text_box(slide8, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "The AI Advisor That Does Everything", font_size=38, color=WHITE, bold=True)

add_text_box(slide8, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "One click. Full analysis. Written recovery plan.",
             font_size=16, color=ACCENT_CYAN)

# Agent flow cards
flow_steps = [
    ("🧑 User", "Enters company\nfinancial data", RGBColor(0x2A, 0x2D, 0x4A), WHITE),
    ("→", "", None, ACCENT_CYAN),
    ("🤖 AI Agent", "Decides which tools\nto call and in what order", RGBColor(0x1A, 0x2A, 0x3A), ACCENT_CYAN),
    ("→", "", None, ACCENT_CYAN),
    ("📊 Tools", "Prediction + SHAP +\nRL Strategy", RGBColor(0x1A, 0x33, 0x2A), ACCENT_GREEN),
    ("→", "", None, ACCENT_CYAN),
    ("📝 Report", "Full recovery plan\nwith actionable steps", RGBColor(0x33, 0x1A, 0x2A), ACCENT_ORANGE),
]

x_pos = Inches(0.8)
for item in flow_steps:
    if item[1] == "" and item[0] == "→":
        # Arrow
        add_text_box(slide8, x_pos, Inches(3.2), Inches(0.5), Inches(1.0),
                     "→", font_size=28, color=item[3], alignment=PP_ALIGN.CENTER)
        x_pos += Inches(0.6)
    else:
        name, desc, bg_color, text_color = item
        card_w = Inches(2.5)
        add_shape_rect(slide8, x_pos, Inches(2.8), card_w, Inches(2.0), bg_color,
                       border_color=RGBColor(0x3A, 0x3D, 0x5A))
        add_text_box(slide8, x_pos + Inches(0.15), Inches(2.95), Inches(2.2), Inches(0.4),
                     name, font_size=16, color=text_color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide8, x_pos + Inches(0.15), Inches(3.5), Inches(2.2), Inches(0.8),
                     desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        x_pos += Inches(2.8)

# Key feature box
add_shape_rect(slide8, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), CARD_BG,
               border_color=RGBColor(0x2A, 0x2D, 0x4A))
add_text_box(slide8, Inches(1.2), Inches(5.45), Inches(11), Inches(0.4),
             "What makes it \"Agentic\"?", font_size=16, color=ACCENT_CYAN, bold=True)

add_multiline_text(slide8, Inches(1.2), Inches(5.9), Inches(11), Inches(0.8), [
    {"text": "✦  The AI makes its own decisions — it chooses when to predict, when to explain, and when to optimize.", "size": 13, "color": LIGHT_GRAY},
    {"text": "✦  Powered by LLaMA 3.3 (70B parameters) — a state-of-the-art open-source language model via Groq.", "size": 13, "color": LIGHT_GRAY},
])

add_slide_number(slide8, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9: Live Demo / Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9)
add_top_accent_line(slide9, ACCENT_BLUE)

add_text_box(slide9, Inches(0.8), Inches(0.4), Inches(6), Inches(0.4),
             "TECH STACK", font_size=12, color=ACCENT_BLUE, bold=True)
add_text_box(slide9, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "How Is It All Built?", font_size=38, color=WHITE, bold=True)

add_text_box(slide9, Inches(0.8), Inches(1.7), Inches(10), Inches(0.6),
             "A clean, modern full-stack architecture — everything runs in the browser.",
             font_size=16, color=LIGHT_GRAY)

# Architecture diagram as styled cards
# Frontend
add_shape_rect(slide9, Inches(0.8), Inches(2.6), Inches(3.5), Inches(2.0), RGBColor(0x1A, 0x2A, 0x3A),
               border_color=ACCENT_BLUE)
add_text_box(slide9, Inches(1.0), Inches(2.75), Inches(3), Inches(0.4),
             "🌐  Frontend", font_size=18, color=ACCENT_BLUE, bold=True)
add_multiline_text(slide9, Inches(1.0), Inches(3.25), Inches(3), Inches(1.2), [
    {"text": "• HTML / CSS / JavaScript", "size": 13, "color": LIGHT_GRAY},
    {"text": "• Chart.js for visualizations", "size": 13, "color": LIGHT_GRAY},
    {"text": "• Glassmorphic UI design", "size": 13, "color": LIGHT_GRAY},
])

# Arrow
add_text_box(slide9, Inches(4.5), Inches(3.2), Inches(0.8), Inches(0.6),
             "⟷", font_size=28, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Backend
add_shape_rect(slide9, Inches(5.3), Inches(2.6), Inches(3.5), Inches(2.0), RGBColor(0x1A, 0x33, 0x2A),
               border_color=ACCENT_GREEN)
add_text_box(slide9, Inches(5.5), Inches(2.75), Inches(3), Inches(0.4),
             "⚙️  Backend", font_size=18, color=ACCENT_GREEN, bold=True)
add_multiline_text(slide9, Inches(5.5), Inches(3.25), Inches(3), Inches(1.2), [
    {"text": "• Python + FastAPI", "size": 13, "color": LIGHT_GRAY},
    {"text": "• REST API endpoints", "size": 13, "color": LIGHT_GRAY},
    {"text": "• Serves everything", "size": 13, "color": LIGHT_GRAY},
])

# Arrow
add_text_box(slide9, Inches(9.0), Inches(3.2), Inches(0.8), Inches(0.6),
             "⟷", font_size=28, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# AI/ML
add_shape_rect(slide9, Inches(9.8), Inches(2.6), Inches(3.0), Inches(2.0), RGBColor(0x2A, 0x1A, 0x33),
               border_color=ACCENT_PURPLE)
add_text_box(slide9, Inches(10.0), Inches(2.75), Inches(2.6), Inches(0.4),
             "🧠  AI / ML", font_size=18, color=ACCENT_PURPLE, bold=True)
add_multiline_text(slide9, Inches(10.0), Inches(3.25), Inches(2.6), Inches(1.2), [
    {"text": "• XGBoost Model", "size": 13, "color": LIGHT_GRAY},
    {"text": "• SHAP Analysis", "size": 13, "color": LIGHT_GRAY},
    {"text": "• PPO RL Agent", "size": 13, "color": LIGHT_GRAY},
])

# LLM card
add_shape_rect(slide9, Inches(5.3), Inches(5.0), Inches(3.5), Inches(1.5), RGBColor(0x1A, 0x2A, 0x3A),
               border_color=ACCENT_CYAN)
add_text_box(slide9, Inches(5.5), Inches(5.15), Inches(3), Inches(0.4),
             "🤖  LLM (AI Brain)", font_size=16, color=ACCENT_CYAN, bold=True)
add_multiline_text(slide9, Inches(5.5), Inches(5.6), Inches(3), Inches(0.8), [
    {"text": "• LLaMA 3.3 70B via Groq", "size": 13, "color": LIGHT_GRAY},
    {"text": "• LangChain orchestration", "size": 13, "color": LIGHT_GRAY},
])

# Open source badge
add_shape_rect(slide9, Inches(0.8), Inches(5.2), Inches(3.5), Inches(1.3), CARD_BG,
               border_color=RGBColor(0x2A, 0x2D, 0x4A))
add_text_box(slide9, Inches(1.0), Inches(5.35), Inches(3.2), Inches(0.3),
             "🔓  100% Open Source Stack", font_size=14, color=ACCENT_GREEN, bold=True)
add_text_box(slide9, Inches(1.0), Inches(5.7), Inches(3.2), Inches(0.6),
             "Every tool and library we used\nis free and open-source.", font_size=12, color=MID_GRAY)

add_slide_number(slide9, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10: Thank You / Q&A
# ══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)
add_top_accent_line(slide10, ACCENT_CYAN)

# Decorative shapes (very dark tinted, near background)
add_shape_rect(slide10, Inches(-1), Inches(4.5), Inches(5), Inches(5),
               RGBColor(0x15, 0x18, 0x30), corner_radius=0.5)

add_shape_rect(slide10, Inches(10), Inches(-1), Inches(5), Inches(5),
               RGBColor(0x18, 0x14, 0x2E), corner_radius=0.5)

# Main content
add_multiline_text(slide10, Inches(2), Inches(1.5), Inches(9.5), Inches(2.5), [
    {"text": "Thank You!", "size": 52, "color": WHITE, "bold": True, "align": PP_ALIGN.CENTER, "spacing": 14},
    {"text": "Questions & Discussion", "size": 24, "color": ACCENT_CYAN, "align": PP_ALIGN.CENTER, "spacing": 20},
])

# Summary card
add_shape_rect(slide10, Inches(2.5), Inches(3.8), Inches(8.5), Inches(2.2), CARD_BG,
               border_color=RGBColor(0x2A, 0x2D, 0x4A))
add_text_box(slide10, Inches(2.8), Inches(3.95), Inches(8), Inches(0.4),
             "What We Built:", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_multiline_text(slide10, Inches(3.5), Inches(4.4), Inches(6.5), Inches(1.4), [
    {"text": "🔮  Predict bankruptcy risk using ML", "size": 14, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER, "spacing": 6},
    {"text": "🔍  Explain the reasons behind the risk", "size": 14, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER, "spacing": 6},
    {"text": "🛡️  Suggest a recovery strategy using RL", "size": 14, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER, "spacing": 6},
    {"text": "🤖  Autonomous AI Advisor ties it all together", "size": 14, "color": LIGHT_GRAY, "align": PP_ALIGN.CENTER},
])

# Logo and credits
if os.path.exists(logo_path):
    slide10.shapes.add_picture(logo_path, Inches(5.9), Inches(6.1), height=Inches(0.5))

add_text_box(slide10, Inches(2), Inches(6.7), Inches(9.5), Inches(0.4),
             "Finance & Economics Club  |  IIT Guwahati  |  Techevince 2026",
             font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide10, 10)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "FEC_Risk_AI_Techevince_Presentation.pptx")
prs.save(output_path)
print(f"\n[OK] Presentation saved to: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
